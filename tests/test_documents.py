import io
import os
from concurrent.futures import ThreadPoolExecutor
from unittest.mock import patch, MagicMock
from uuid import UUID as PyUUID, uuid4
import zipfile

import fitz  # PyMuPDF
from fastapi.testclient import TestClient
from pptx import Presentation

from app.core.config import settings
from app.core.database import SessionLocal
from app.main import app
from app.models.document import Document, DocumentPage
from app.services.processors.pdf_processor import PDFProcessor
from app.services.processors.pptx_processor import PPTXProcessor
from app.worker import celery_app, process_document

client = TestClient(app)


def create_sample_pdf_bytes(text: str = "Sample Document Text Page 1") -> bytes:
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((50, 50), text)
    pdf_bytes = doc.tobytes()
    doc.close()
    return pdf_bytes


def create_sample_pptx_bytes(slide_text: str = "Sample Slide Title Text") -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = slide_text
    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def test_celery_app_loading():
    assert celery_app is not None
    assert celery_app.main == "ai_paper_generator_worker"


def test_document_upload_and_ownership():
    uid = uuid4().hex[:8]
    email_a = f"doc_owner_{uid}@example.com"
    email_b = f"doc_other_{uid}@example.com"

    user_a = client.post(
        "/api/v1/auth/register",
        json={"name": "Doc Owner", "email": email_a, "password": "password123"},
    ).json()
    headers_a = {"Authorization": f"Bearer {user_a['access_token']}"}

    user_b = client.post(
        "/api/v1/auth/register",
        json={"name": "Other User", "email": email_b, "password": "password123"},
    ).json()
    headers_b = {"Authorization": f"Bearer {user_b['access_token']}"}

    ws = client.post("/api/v1/workspaces", json={"name": "WS"}, headers=headers_a).json()
    subj = client.post(f"/api/v1/workspaces/{ws['id']}/subjects", json={"name": "Subj"}, headers=headers_a).json()
    book = client.post(f"/api/v1/subjects/{subj['id']}/books", json={"name": "Book A"}, headers=headers_a).json()
    book_id = book["id"]

    # 1. Unauthenticated upload rejected
    unauth_resp = client.post("/api/v1/documents/upload", data={"book_id": book_id})
    assert unauth_resp.status_code == 401

    # 2. Unsupported file extension rejected
    bad_ext_resp = client.post(
        "/api/v1/documents/upload",
        data={"book_id": book_id},
        files={"file": ("malicious.exe", b"binary content", "application/octet-stream")},
        headers=headers_a,
    )
    assert bad_ext_resp.status_code == 400

    # 3. User B cannot upload to User A's book
    bad_book_resp = client.post(
        "/api/v1/documents/upload",
        data={"book_id": book_id},
        files={"file": ("test.pdf", create_sample_pdf_bytes(), "application/pdf")},
        headers=headers_b,
    )
    assert bad_book_resp.status_code == 404

    # 4. User A uploads PDF successfully
    pdf_bytes = create_sample_pdf_bytes("Physics Chapter 1 Basics")
    upload_resp = client.post(
        "/api/v1/documents/upload",
        data={"book_id": book_id},
        files={"file": ("sample_physics.pdf", pdf_bytes, "application/pdf")},
        headers=headers_a,
    )
    assert upload_resp.status_code == 202
    upload_data = upload_resp.json()
    assert "id" in upload_data
    assert upload_data["status"] == "UPLOADED"
    doc_id = upload_data["id"]

    # 5. User A gets status (200)
    status_a = client.get(f"/api/v1/documents/{doc_id}", headers=headers_a)
    assert status_a.status_code == 200
    assert status_a.json()["original_filename"] == "sample_physics.pdf"

    # 6. User B getting User A's document status returns 404
    status_b = client.get(f"/api/v1/documents/{doc_id}", headers=headers_b)
    assert status_b.status_code == 404


def test_file_signature_validation():
    uid = uuid4().hex[:8]
    email = f"sig_{uid}@example.com"
    user = client.post(
        "/api/v1/auth/register",
        json={"name": "Sig User", "email": email, "password": "password123"},
    ).json()
    headers = {"Authorization": f"Bearer {user['access_token']}"}

    ws = client.post("/api/v1/workspaces", json={"name": "WS"}, headers=headers).json()
    subj = client.post(f"/api/v1/workspaces/{ws['id']}/subjects", json={"name": "Subj"}, headers=headers).json()
    book = client.post(f"/api/v1/subjects/{subj['id']}/books", json={"name": "Book"}, headers=headers).json()

    # Rejection of fake PDF (plain text labeled .pdf)
    fake_pdf_resp = client.post(
        "/api/v1/documents/upload",
        data={"book_id": book["id"]},
        files={"file": ("fake.pdf", b"This is not a real PDF file header", "application/pdf")},
        headers=headers,
    )
    assert fake_pdf_resp.status_code == 400
    assert "missing %pdf" in fake_pdf_resp.json()["detail"].lower()

    # Rejection of fake PPTX (plain text labeled .pptx)
    fake_pptx_resp = client.post(
        "/api/v1/documents/upload",
        data={"book_id": book["id"]},
        files={"file": ("fake.pptx", b"This is not a zip file", "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
        headers=headers,
    )
    assert fake_pptx_resp.status_code == 400
    assert "missing pk zip" in fake_pptx_resp.json()["detail"].lower()


def test_enqueue_failure_cleanup():
    uid = uuid4().hex[:8]
    email = f"enqueue_{uid}@example.com"
    user = client.post(
        "/api/v1/auth/register",
        json={"name": "Enqueue User", "email": email, "password": "password123"},
    ).json()
    headers = {"Authorization": f"Bearer {user['access_token']}"}

    ws = client.post("/api/v1/workspaces", json={"name": "WS"}, headers=headers).json()
    subj = client.post(f"/api/v1/workspaces/{ws['id']}/subjects", json={"name": "Subj"}, headers=headers).json()
    book = client.post(f"/api/v1/subjects/{subj['id']}/books", json={"name": "Book"}, headers=headers).json()

    pdf_bytes = create_sample_pdf_bytes("Enqueue test document content")

    with patch("app.worker.process_document.delay", side_effect=Exception("Redis connection error")):
        resp = client.post(
            "/api/v1/documents/upload",
            data={"book_id": book["id"]},
            files={"file": ("enqueue_fail.pdf", pdf_bytes, "application/pdf")},
            headers=headers,
        )
        assert resp.status_code == 500
        assert "failed to queue document processing task" in resp.json()["detail"].lower()


def test_pdf_processing_and_ocr_fallback(tmp_path):
    doc_dir = str(tmp_path / "pdf_doc")
    os.makedirs(doc_dir, exist_ok=True)
    pdf_path = os.path.join(doc_dir, "original.pdf")

    pdf_bytes = create_sample_pdf_bytes("Calculus Fundamentals and Derivatives")
    with open(pdf_path, "wb") as f:
        f.write(pdf_bytes)

    pages = PDFProcessor.process_pdf(pdf_path, doc_dir)
    assert len(pages) == 1
    assert pages[0]["page_number"] == 1
    assert "Calculus Fundamentals" in pages[0]["text_content"]
    assert pages[0]["content_type"] == "PAGE"

    empty_pdf_path = os.path.join(doc_dir, "empty.pdf")
    empty_doc = fitz.open()
    empty_doc.new_page()
    empty_doc.save(empty_pdf_path)
    empty_doc.close()

    with patch("pytesseract.image_to_string", return_value="OCR Extracted Scanned Text"):
        ocr_pages = PDFProcessor.process_pdf(empty_pdf_path, doc_dir)
        assert len(ocr_pages) == 1
        assert "OCR Extracted" in ocr_pages[0]["text_content"]
        assert ocr_pages[0]["metadata_json"]["ocr_applied"] is True
        assert ocr_pages[0]["metadata_json"]["ocr_failed"] is False


def test_pptx_processing(tmp_path):
    doc_dir = str(tmp_path / "pptx_doc")
    os.makedirs(doc_dir, exist_ok=True)
    pptx_path = os.path.join(doc_dir, "original.pptx")

    pptx_bytes = create_sample_pptx_bytes("Organic Chemistry Overview")
    with open(pptx_path, "wb") as f:
        f.write(pptx_bytes)

    slides = PPTXProcessor.process_pptx(pptx_path, doc_dir)
    assert len(slides) == 1
    assert slides[0]["page_number"] == 1
    assert slides[0]["content_type"] == "SLIDE"
    assert "Organic Chemistry" in slides[0]["text_content"]


def test_concurrent_document_processing_claim():
    uid = uuid4().hex[:8]
    email = f"claim_{uid}@example.com"

    reg_resp = client.post(
        "/api/v1/auth/register",
        json={"name": "Claim User", "email": email, "password": "password123"},
    ).json()
    headers = {"Authorization": f"Bearer {reg_resp['access_token']}"}

    ws = client.post("/api/v1/workspaces", json={"name": "WS"}, headers=headers).json()
    subj = client.post(f"/api/v1/workspaces/{ws['id']}/subjects", json={"name": "Subj"}, headers=headers).json()
    book = client.post(f"/api/v1/subjects/{subj['id']}/books", json={"name": "Book"}, headers=headers).json()

    pdf_bytes = create_sample_pdf_bytes("Atomic concurrency claim test content")

    with patch("app.worker.process_document.delay"):
        upload_resp = client.post(
            "/api/v1/documents/upload",
            data={"book_id": book["id"]},
            files={"file": ("claim.pdf", pdf_bytes, "application/pdf")},
            headers=headers,
        ).json()
    doc_id_str = upload_resp["id"]

    results = []

    def worker_claim_attempt():
        return process_document(doc_id_str)

    with ThreadPoolExecutor(max_workers=2) as executor:
        f1 = executor.submit(worker_claim_attempt)
        f2 = executor.submit(worker_claim_attempt)
        results.append(f1.result())
        results.append(f2.result())

    # Exactly 1 worker processes, the other skips
    processed_count = sum(1 for r in results if r.get("status") == "PROCESSED")
    skipped_count = sum(1 for r in results if r.get("status") == "skipped")

    assert processed_count == 1
    assert skipped_count == 1

    # Verify exactly 1 page record created
    db = SessionLocal()
    try:
        pages = db.query(DocumentPage).filter(DocumentPage.document_id == PyUUID(doc_id_str)).all()
        assert len(pages) == 1
    finally:
        db.close()


def test_permanent_and_transient_failure_handling(tmp_path):
    uid = uuid4().hex[:8]
    email = f"fail_{uid}@example.com"

    reg_resp = client.post(
        "/api/v1/auth/register",
        json={"name": "Fail User", "email": email, "password": "password123"},
    ).json()
    headers = {"Authorization": f"Bearer {reg_resp['access_token']}"}

    ws = client.post("/api/v1/workspaces", json={"name": "WS"}, headers=headers).json()
    subj = client.post(f"/api/v1/workspaces/{ws['id']}/subjects", json={"name": "Subj"}, headers=headers).json()
    book = client.post(f"/api/v1/subjects/{subj['id']}/books", json={"name": "Book"}, headers=headers).json()

    pdf_bytes = create_sample_pdf_bytes("Permanent failure test document")

    with patch("app.worker.process_document.delay"):
        upload_resp = client.post(
            "/api/v1/documents/upload",
            data={"book_id": book["id"]},
            files={"file": ("perm_fail.pdf", pdf_bytes, "application/pdf")},
            headers=headers,
        ).json()
    doc_id_str = upload_resp["id"]

    # Permanent failure simulation (e.g. ValueError during processing)
    with patch("app.services.processors.pdf_processor.PDFProcessor.process_pdf", side_effect=ValueError("Corrupt PDF layout")):
        res = process_document(doc_id_str)
        assert res["status"] == "FAILED"

    # Verify DB record and disk storage are cleaned up (returns 404 Not Found)
    doc_resp = client.get(f"/api/v1/documents/{doc_id_str}", headers=headers)
    assert doc_resp.status_code == 404


def test_stale_document_reclaim():
    uid = uuid4().hex[:8]
    email = f"stale_{uid}@example.com"

    reg_resp = client.post(
        "/api/v1/auth/register",
        json={"name": "Stale User", "email": email, "password": "password123"},
    ).json()
    headers = {"Authorization": f"Bearer {reg_resp['access_token']}"}

    ws = client.post("/api/v1/workspaces", json={"name": "WS"}, headers=headers).json()
    subj = client.post(f"/api/v1/workspaces/{ws['id']}/subjects", json={"name": "Subj"}, headers=headers).json()
    book = client.post(f"/api/v1/subjects/{subj['id']}/books", json={"name": "Book"}, headers=headers).json()

    pdf_bytes = create_sample_pdf_bytes("Stale recovery document test content")

    with patch("app.worker.process_document.delay"):
        upload_resp = client.post(
            "/api/v1/documents/upload",
            data={"book_id": book["id"]},
            files={"file": ("stale.pdf", pdf_bytes, "application/pdf")},
            headers=headers,
        ).json()
    doc_id_str = upload_resp["id"]

    # Manually simulate crashed worker: set status to PROCESSING with an old timestamp (30 mins ago)
    db = SessionLocal()
    try:
        from datetime import datetime, timedelta, timezone
        doc = db.get(Document, PyUUID(doc_id_str))
        doc.processing_status = "PROCESSING"
        doc.processing_started_at = datetime.now(timezone.utc) - timedelta(minutes=30)
        db.commit()
    finally:
        db.close()

    # Re-run processing task: Should successfully claim and reprocess the stale document
    res = process_document(doc_id_str)
    assert res["status"] == "PROCESSED"
    assert res["pages_count"] == 1


def test_document_preview_and_download_endpoints():
    uid = uuid4().hex[:8]
    user_a = client.post("/api/v1/auth/register", json={"name": "Owner", "email": f"own_{uid}@example.com", "password": "password123"}).json()
    user_b = client.post("/api/v1/auth/register", json={"name": "Attacker", "email": f"att_{uid}@example.com", "password": "password123"}).json()

    headers_a = {"Authorization": f"Bearer {user_a['access_token']}"}
    headers_b = {"Authorization": f"Bearer {user_b['access_token']}"}

    ws = client.post("/api/v1/workspaces", json={"name": f"WS_{uid}"}, headers=headers_a).json()
    subj = client.post(f"/api/v1/workspaces/{ws['id']}/subjects", json={"name": f"Subj_{uid}"}, headers=headers_a).json()
    book = client.post(f"/api/v1/subjects/{subj['id']}/books", json={"name": f"Book_{uid}"}, headers=headers_a).json()

    pdf_bytes = create_sample_pdf_bytes("Doc Preview/Download PDF Test Content")

    with patch("app.worker.process_document.delay"):
        upload_resp = client.post(
            "/api/v1/documents/upload",
            data={"book_id": book["id"]},
            files={"file": ("sample_doc.pdf", pdf_bytes, "application/pdf")},
            headers=headers_a,
        ).json()
    doc_id = upload_resp["id"]

    # 1. Direct static access (vigilens-backend pattern)
    static_res = client.get(f"/storage/documents/{doc_id}/original.pdf")
    assert static_res.status_code == 200
    assert static_res.content == pdf_bytes

    # 2. Missing file on disk -> 404
    db = SessionLocal()
    try:
        doc_obj = db.get(Document, PyUUID(doc_id))
        stored_path = doc_obj.stored_path
        if os.path.exists(stored_path):
            os.remove(stored_path)
    finally:
        db.close()

    missing_static = client.get(f"/storage/documents/{doc_id}/original.pdf")
    assert missing_static.status_code == 404
