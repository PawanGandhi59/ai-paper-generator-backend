import os
from typing import Any, Dict, List

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.exc import PackageNotFoundError


class PPTXProcessor:
    @staticmethod
    def process_pptx(file_path: str, doc_dir: str) -> List[Dict[str, Any]]:
        slides_data = []

        if not os.path.exists(file_path):
            raise FileNotFoundError(f"PPTX file not found at path: {file_path}")

        try:
            prs = Presentation(file_path)
        except (PackageNotFoundError, Exception) as open_exc:
            raise ValueError(f"Corrupt or invalid PPTX file: {open_exc}")

        for slide_index, slide in enumerate(prs.slides):
            slide_number = slide_index + 1
            extracted_text_chunks = []

            slide_images_dir = os.path.join(doc_dir, "pages", f"{slide_number:04d}", "images")
            os.makedirs(slide_images_dir, exist_ok=True)
            saved_image_paths = []
            img_counter = 0

            for shape in slide.shapes:
                # 1. Text extraction from text frames
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        p_text = paragraph.text.strip()
                        if p_text:
                            extracted_text_chunks.append(p_text)

                # 2. Text extraction from tables
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if row_text:
                            extracted_text_chunks.append(row_text)

                # 3. Image extraction from picture shapes
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or hasattr(shape, "image"):
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        image_ext = image.ext

                        img_counter += 1
                        image_filename = f"image_{img_counter}.{image_ext}"
                        full_img_path = os.path.join(slide_images_dir, image_filename)

                        with open(full_img_path, "wb") as img_file:
                            img_file.write(image_bytes)

                        saved_image_paths.append(full_img_path)
                    except Exception as img_exc:
                        print(f"PPTX image extraction warning on slide {slide_number}: {img_exc}")

            slide_text = "\n".join(extracted_text_chunks)
            first_image_path = saved_image_paths[0] if saved_image_paths else None

            slides_data.append({
                "page_number": slide_number,
                "content_type": "SLIDE",
                "text_content": slide_text,
                "image_path": first_image_path,
                "metadata_json": {
                    "shape_count": len(slide.shapes),
                    "image_count": len(saved_image_paths),
                    "extracted_image_paths": saved_image_paths,
                },
            })

        return slides_data
